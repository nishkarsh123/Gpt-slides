
from datetime import datetime
from celery import shared_task
from .utils import send_generation_request_to_gpt
from .models import PPT
from pptx import Presentation	 
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rest_framework.response import Response
from django.conf import settings
import logging
import os
 
logger = logging.getLogger('opai')
MEDIA_ROOT=settings.MEDIA_ROOT

def find_shape_by_text(shapes, text):
    for shape in shapes:
        if shape.has_text_frame:
            if shape.text.lower() == text.lower():
                return shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return find_shape_by_text(shape.shapes, text)
    return None

@shared_task
def send_generation_request_to_GPT_task(id):
    try:
        # unique_file_path = MEDIA_ROOT / "ppt" / f"{unique_filename}"
        # new_ppt = PPT.objects.filter(id=id).first()

        instance = PPT.objects.get(id=id)
        result = send_generation_request_to_gpt()
        if result:
            root = Presentation(instance.ppt.path)
            slide = root.slides[0]
            title = find_shape_by_text(slide.shapes, "THOUGHT/JOKE/QUOTE OF THE DAY")
            subTitle = find_shape_by_text(slide.shapes, "{DYNAMIC_FIELD FOR THOUGHT/JOKE/QUOTE}")

            if title:
                title.text = result['title']
            else:
                slide.shapes.title.text = result['title']

            if subTitle:
                subTitle.text = result['joke']
            elif slide.placeholders[1]:
                slide.placeholders[1].text = result['joke']

            new_file_path = f"ppt_modified/{os.path.basename(instance.ppt.name)}"
            unique_file_path = MEDIA_ROOT / "ppt_modified" / f"{os.path.basename(instance.ppt.name)}"
            root.save(unique_file_path)

            instance.ppt_modified = new_file_path
            instance.status = "finished"
            instance.save()

            return {"id": instance.id, "status": instance.status, "created_at": instance.created_at.isoformat()}
        else:
            raise Exception("Result is empty")

    except ConnectionError as ce:
        logger.error(f"Connection error: {ce}")
        
    except Exception as e:
        logger.exception("An error occurred: %s", e)
        instance.status = "failed"
        instance.save()
        raise Exception(f"Error: {e}") from None